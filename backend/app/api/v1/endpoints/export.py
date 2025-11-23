from fastapi import APIRouter, Depends, HTTPException, status, Query
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from uuid import UUID
import os
from pathlib import Path
from app.core.database import get_db
from app.api.v1.dependencies import get_current_user
from app.models.user import User
from app.models.project import Project, DocumentType
from app.models.section import Section
from app.core.config import settings
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

router = APIRouter()

def ensure_export_dir():
    """Ensure export directory exists."""
    export_dir = Path(settings.EXPORT_TMP_DIR)
    export_dir.mkdir(parents=True, exist_ok=True)
    return export_dir

# Ensure export directory exists on startup
ensure_export_dir()

@router.get("/project/{project_id}")
async def export_project(
    project_id: UUID,
    type: str = Query(..., description="Export type: docx or pptx"),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Export project as .docx or .pptx file."""
    # Verify project belongs to user
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Project not found"
        )
    
    # Get all sections ordered by order_index
    sections = db.query(Section).filter(
        Section.project_id == project_id
    ).order_by(Section.order_index).all()
    
    if not sections:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Project has no sections to export"
        )
    
    export_dir = ensure_export_dir()
    
    if type.lower() == "docx":
        if project.doc_type != DocumentType.DOCX:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Project is not a Word document"
            )
        
        # Create Word document
        doc = Document()
        doc.add_heading(project.title, 0)
        
        for section in sections:
            # Add section heading
            doc.add_heading(section.title, level=1)
            
            # Add content
            if section.content:
                # Split content into paragraphs
                paragraphs = section.content.split('\n\n')
                for para in paragraphs:
                    if para.strip():
                        doc.add_paragraph(para.strip())
            else:
                doc.add_paragraph("[Content not generated]")
        
        # Save file
        filename = f"{project.title.replace(' ', '_')}_{project_id}.docx"
        filepath = export_dir / filename
        doc.save(str(filepath))
        
        return FileResponse(
            path=str(filepath),
            filename=filename,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
    
    elif type.lower() == "pptx":
        if project.doc_type != DocumentType.PPTX:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Project is not a PowerPoint presentation"
            )
        
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = project.title
        subtitle.text = project.topic
        
        # Content slides
        for section in sections:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = section.title
            
            tf = body_shape.text_frame
            tf.text = section.content or "[Content not generated]"
            
            # Format text
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(14)
        
        # Save file
        filename = f"{project.title.replace(' ', '_')}_{project_id}.pptx"
        filepath = export_dir / filename
        prs.save(str(filepath))
        
        return FileResponse(
            path=str(filepath),
            filename=filename,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    
    else:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Invalid export type. Use 'docx' or 'pptx'"
        )

